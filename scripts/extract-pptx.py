#!/usr/bin/env python3
"""
Extract all content from a PowerPoint file (.pptx).
Returns a JSON structure with slides, text, and images.

Usage:
    python extract-pptx.py <input.pptx> [output_dir]

Requires: pip install python-pptx
"""

import json
import os
import sys

# ── Pre-flight: check dependency ─────────────────────────────
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("错误: 缺少依赖 python-pptx", file=sys.stderr)
    print("", file=sys.stderr)
    print("请安装后重试:", file=sys.stderr)
    print("  pip install python-pptx", file=sys.stderr)
    print("", file=sys.stderr)
    print("如果是 Windows 权限问题，请尝试:", file=sys.stderr)
    print("  pip install --user python-pptx", file=sys.stderr)
    sys.exit(2)

def extract_pptx(file_path, output_dir="."):
    """
    Extract all content from a PowerPoint file.
    Returns a list of slide data dicts with text, images, and notes.
    """
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"错误: 无法打开 PPT 文件 '{file_path}': {e}", file=sys.stderr)
        sys.exit(3)

    slides_data = []

    # Create assets directory for extracted images
    assets_dir = os.path.join(output_dir, "assets")
    os.makedirs(assets_dir, exist_ok=True)

    for slide_num, slide in enumerate(prs.slides):
        slide_data = {
            "number": slide_num + 1,
            "title": "",
            "content": [],
            "images": [],
            "notes": "",
        }

        try:
            for shape in slide.shapes:
                # Extract text content
                if shape.has_text_frame:
                    if shape == slide.shapes.title:
                        slide_data["title"] = shape.text
                    else:
                        slide_data["content"].append(
                            {"type": "text", "content": shape.text}
                        )

                # Extract images (shape_type 13 = Picture)
                if hasattr(shape, 'image') and shape.shape_type == 13:
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        image_ext = image.ext
                        image_name = f"slide{slide_num + 1}_img{len(slide_data['images']) + 1}.{image_ext}"
                        image_path = os.path.join(assets_dir, image_name)

                        with open(image_path, "wb") as f:
                            f.write(image_bytes)

                        slide_data["images"].append(
                            {
                                "path": f"assets/{image_name}",
                                "width": shape.width,
                                "height": shape.height,
                            }
                        )
                    except Exception as e:
                        print(f"警告: slide {slide_num+1} 的图片提取失败: {e}", file=sys.stderr)

            # Extract speaker notes
            if slide.has_notes_slide:
                try:
                    notes_frame = slide.notes_slide.notes_text_frame
                    slide_data["notes"] = notes_frame.text
                except Exception:
                    pass
        except Exception as e:
            print(f"警告: slide {slide_num+1} 处理失败 (已跳过): {e}", file=sys.stderr)
            continue

        slides_data.append(slide_data)

    return slides_data

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extract-pptx.py <input.pptx> [output_dir]")
        sys.exit(1)

    input_file = sys.argv[1]
    if not os.path.exists(input_file):
        print(f"错误: 文件不存在: {input_file}", file=sys.stderr)
        sys.exit(3)

    output_dir = sys.argv[2] if len(sys.argv) > 2 else "."

    try:
        slides = extract_pptx(input_file, output_dir)
    except Exception as e:
        print(f"错误: 提取过程出现异常: {e}", file=sys.stderr)
        sys.exit(4)

    # Write extracted data as JSON
    output_path = os.path.join(output_dir, "extracted-slides.json")
    try:
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(slides, f, indent=2, ensure_ascii=False)
    except Exception as e:
        print(f"错误: 无法写入输出文件 '{output_path}': {e}", file=sys.stderr)
        sys.exit(4)

    print(f"提取完成: {len(slides)} 张幻灯片 → {output_path}")
    for s in slides:
        img_count = len(s["images"])
        print(f"  Slide {s['number']}: {s['title'] or '(无标题)'} — {img_count} 张图片")
